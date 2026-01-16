"""PPTX text extraction."""

from typing import List, Dict
import logging

logger = logging.getLogger(__name__)


def extract_text_from_pptx(file_path: str) -> List[Dict]:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        pages = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            # Combine all text from the slide
            slide_text = "\n\n".join(slide_texts)
            
            if slide_text:
                pages.append({
                    "page_number": slide_num,
                    "text": slide_text
                })
        
        return pages if pages else [{"page_number": 1, "text": ""}]
        
    except ImportError:
        logger.error("python-pptx not installed. Install it with: pip install python-pptx")
        raise Exception("PPTX processing requires python-pptx package")
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")

