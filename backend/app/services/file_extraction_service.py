"""File extraction service for multiple file formats."""

import mimetypes
import os
from typing import List, Dict, Optional
import logging

logger = logging.getLogger(__name__)


def get_file_type(filename: str) -> tuple[str, str]:
    """
    Determine file type and MIME type from filename.
    
    Args:
        filename: Name of the file
        
    Returns:
        Tuple of (file_type, mime_type)
    """
    ext = os.path.splitext(filename)[1].lower()
    
    # Map extensions to file types and MIME types
    type_map = {
        '.pdf': ('pdf', 'application/pdf'),
        '.docx': ('docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
        '.doc': ('doc', 'application/msword'),
        '.pptx': ('pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
        '.ppt': ('ppt', 'application/vnd.ms-powerpoint'),
        '.txt': ('txt', 'text/plain'),
        '.rtf': ('rtf', 'application/rtf'),
        '.csv': ('csv', 'text/csv'),
        '.xlsx': ('xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'),
        '.xls': ('xls', 'application/vnd.ms-excel'),
        '.mp3': ('mp3', 'audio/mpeg'),
        '.wav': ('wav', 'audio/wav'),
        '.mp4': ('mp4', 'video/mp4'),
        '.webm': ('webm', 'video/webm'),
        '.png': ('png', 'image/png'),
        '.jpg': ('jpg', 'image/jpeg'),
        '.jpeg': ('jpeg', 'image/jpeg'),
        '.svg': ('svg', 'image/svg+xml'),
    }
    
    file_type, mime_type = type_map.get(ext, ('unknown', 'application/octet-stream'))
    
    # Fallback to mimetypes if not in our map
    if file_type == 'unknown':
        guessed_mime, _ = mimetypes.guess_type(filename)
        if guessed_mime:
            mime_type = guessed_mime
            # Extract file type from MIME type
            if '/' in guessed_mime:
                file_type = guessed_mime.split('/')[1].split('-')[-1]
    
    return file_type, mime_type


def is_text_processable(file_type: str) -> bool:
    """
    Check if a file type should be processed for text extraction and RAG.
    
    Args:
        file_type: File type string
        
    Returns:
        True if file should be processed for text extraction
    """
    text_processable = {'pdf', 'docx', 'doc', 'pptx', 'ppt', 'txt', 'rtf', 'csv', 'xlsx', 'xls'}
    return file_type in text_processable


def is_media_file(file_type: str) -> bool:
    """
    Check if a file is a media file (audio/video).
    
    Args:
        file_type: File type string
        
    Returns:
        True if file is a media file
    """
    media_types = {'mp3', 'wav', 'mp4', 'webm'}
    return file_type in media_types


def is_image_file(file_type: str) -> bool:
    """
    Check if a file is an image file.
    
    Args:
        file_type: File type string
        
    Returns:
        True if file is an image file
    """
    image_types = {'png', 'jpg', 'jpeg', 'svg'}
    return file_type in image_types


def extract_text_from_file(file_path: str, file_type: str) -> List[Dict]:
    """
    Extract text from various file formats.
    
    Args:
        file_path: Path to the file
        file_type: Type of the file (pdf, docx, pptx, etc.)
        
    Returns:
        List of dictionaries with page_number and text keys
        
    Raises:
        Exception: If file cannot be processed
    """
    if file_type == 'pdf':
        from app.services.pdf_service import extract_text_from_pdf
        return extract_text_from_pdf(file_path)
    elif file_type in ('docx', 'doc'):
        return extract_text_from_docx(file_path)
    elif file_type in ('pptx', 'ppt'):
        return extract_text_from_pptx(file_path)
    elif file_type == 'txt':
        return extract_text_from_txt(file_path)
    elif file_type == 'rtf':
        return extract_text_from_rtf(file_path)
    elif file_type == 'csv':
        return extract_text_from_csv(file_path)
    elif file_type in ('xlsx', 'xls'):
        return extract_text_from_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type for text extraction: {file_type}")


def extract_text_from_docx(file_path: str) -> List[Dict]:
    """Extract text from DOCX file."""
    try:
        from docx import Document as DocxDocument
        
        doc = DocxDocument(file_path)
        pages = []
        
        # DOCX doesn't have pages, so we'll treat paragraphs as "pages"
        # Group paragraphs into logical pages (every ~500 words)
        current_page_text = []
        current_page_num = 1
        word_count = 0
        words_per_page = 500
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                words = text.split()
                word_count += len(words)
                current_page_text.append(text)
                
                # If we've accumulated enough words, create a page
                if word_count >= words_per_page:
                    pages.append({
                        "page_number": current_page_num,
                        "text": "\n\n".join(current_page_text)
                    })
                    current_page_text = []
                    current_page_num += 1
                    word_count = 0
        
        # Add remaining text as last page
        if current_page_text:
            pages.append({
                "page_number": current_page_num,
                "text": "\n\n".join(current_page_text)
            })
        
        # If no pages were created, create one with all text
        if not pages:
            all_text = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if all_text:
                pages.append({
                    "page_number": 1,
                    "text": all_text
                })
        
        return pages if pages else [{"page_number": 1, "text": ""}]
        
    except ImportError:
        logger.error("python-docx not installed. Install it with: pip install python-docx")
        raise Exception("DOCX processing requires python-docx package")
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_pptx(file_path: str) -> List[Dict]:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        pages = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            # Combine all text from the slide
            slide_text = "\n\n".join(slide_texts)
            
            if slide_text:
                pages.append({
                    "page_number": slide_num,
                    "text": slide_text
                })
        
        return pages if pages else [{"page_number": 1, "text": ""}]
        
    except ImportError:
        logger.error("python-pptx not installed. Install it with: pip install python-pptx")
        raise Exception("PPTX processing requires python-pptx package")
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")


def extract_text_from_txt(file_path: str) -> List[Dict]:
    """Extract text from TXT file."""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        text = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if text is None:
            # Last resort: read as binary and decode with errors='replace'
            with open(file_path, 'rb') as f:
                text = f.read().decode('utf-8', errors='replace')
        
        # Split into pages (every ~2000 characters)
        chars_per_page = 2000
        pages = []
        page_num = 1
        
        for i in range(0, len(text), chars_per_page):
            page_text = text[i:i + chars_per_page]
            pages.append({
                "page_number": page_num,
                "text": page_text
            })
            page_num += 1
        
        return pages if pages else [{"page_number": 1, "text": text or ""}]
        
    except Exception as e:
        logger.error(f"Error extracting text from TXT {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from TXT: {str(e)}")


def extract_text_from_rtf(file_path: str) -> List[Dict]:
    """Extract text from RTF file."""
    try:
        import re
        
        with open(file_path, 'r', encoding='latin-1', errors='ignore') as f:
            rtf_content = f.read()
        
        # Simple RTF text extraction (removes RTF control codes)
        # This is a basic implementation - for production, consider using striprtf or similar
        text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
        text = re.sub(r'\{[^}]*\}', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        
        # Split into pages
        chars_per_page = 2000
        pages = []
        page_num = 1
        
        for i in range(0, len(text), chars_per_page):
            page_text = text[i:i + chars_per_page]
            pages.append({
                "page_number": page_num,
                "text": page_text
            })
            page_num += 1
        
        return pages if pages else [{"page_number": 1, "text": text or ""}]
        
    except Exception as e:
        logger.error(f"Error extracting text from RTF {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from RTF: {str(e)}")


def extract_text_from_csv(file_path: str) -> List[Dict]:
    """Extract text from CSV file."""
    try:
        import csv
        
        pages = []
        page_num = 1
        rows_per_page = 100
        current_page_rows = []
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            # Try to detect delimiter
            sample = f.read(1024)
            f.seek(0)
            sniffer = csv.Sniffer()
            delimiter = sniffer.sniff(sample).delimiter
            
            reader = csv.reader(f, delimiter=delimiter)
            
            # Read header
            try:
                header = next(reader)
                current_page_rows.append(" | ".join(header))
            except StopIteration:
                pass
            
            # Read rows
            for row in reader:
                current_page_rows.append(" | ".join(row))
                
                if len(current_page_rows) >= rows_per_page:
                    pages.append({
                        "page_number": page_num,
                        "text": "\n".join(current_page_rows)
                    })
                    current_page_rows = []
                    page_num += 1
            
            # Add remaining rows
            if current_page_rows:
                pages.append({
                    "page_number": page_num,
                    "text": "\n".join(current_page_rows)
                })
        
        return pages if pages else [{"page_number": 1, "text": ""}]
        
    except Exception as e:
        logger.error(f"Error extracting text from CSV {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from CSV: {str(e)}")


def extract_text_from_xlsx(file_path: str) -> List[Dict]:
    """Extract text from XLSX file."""
    try:
        try:
            import openpyxl
            use_openpyxl = True
        except ImportError:
            try:
                import pandas as pd
                use_openpyxl = False
            except ImportError:
                raise ImportError("Neither openpyxl nor pandas is installed. Install one with: pip install openpyxl or pip install pandas")
        
        pages = []
        
        if use_openpyxl:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            page_num = 1
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_texts = []
                
                # Add sheet name as header
                sheet_texts.append(f"Sheet: {sheet_name}")
                sheet_texts.append("")
                
                # Extract data from cells
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_texts.append(row_text)
                
                if sheet_texts:
                    pages.append({
                        "page_number": page_num,
                        "text": "\n".join(sheet_texts)
                    })
                    page_num += 1
            
            wb.close()
        else:
            # Use pandas
            df_dict = pd.read_excel(file_path, sheet_name=None)
            page_num = 1
            
            for sheet_name, df in df_dict.items():
                sheet_texts = [f"Sheet: {sheet_name}", ""]
                sheet_texts.append(df.to_string())
                
                pages.append({
                    "page_number": page_num,
                    "text": "\n".join(sheet_texts)
                })
                page_num += 1
        
        return pages if pages else [{"page_number": 1, "text": ""}]
        
    except Exception as e:
        logger.error(f"Error extracting text from XLSX {file_path}: {str(e)}")
        raise Exception(f"Failed to extract text from XLSX: {str(e)}")

