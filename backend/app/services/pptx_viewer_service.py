"""PPTX slide extraction service for viewing."""

import logging
from typing import List, Dict

logger = logging.getLogger(__name__)


def extract_pptx_slides(file_path: str) -> List[Dict]:
    """
    Extract slides from PPTX file for viewing.
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        List of slide dictionaries with content
        
    Raises:
        Exception: If extraction fails
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = {
                "slide_number": slide_num,
                "title": "",
                "text_content": [],
                "shapes": []
            }
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                shape_info = {
                    "type": shape.shape_type,
                    "text": ""
                }
                
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    shape_info["text"] = text
                    slide_content["text_content"].append(text)
                    
                    # Check if it's a title (usually first text box or placeholder)
                    if not slide_content["title"] and len(text) < 200:
                        slide_content["title"] = text
                
                slide_content["shapes"].append(shape_info)
            
            # If no title found, use first text content
            if not slide_content["title"] and slide_content["text_content"]:
                slide_content["title"] = slide_content["text_content"][0][:100]
            
            # Combine all text content
            slide_content["full_text"] = "\n\n".join(slide_content["text_content"])
            
            slides.append(slide_content)
        
        return slides
        
    except ImportError:
        logger.error("python-pptx not installed")
        raise Exception("PPTX processing requires python-pptx package")
    except Exception as e:
        logger.error(f"Error extracting PPTX slides {file_path}: {str(e)}")
        raise Exception(f"Failed to extract PPTX slides: {str(e)}")

